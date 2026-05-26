from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
from PIL import Image

OUT_DIR = 'presentations'
OUT_FILE = os.path.join(OUT_DIR, 'SmartFare_Tracker.pptx')
SCREENSHOT = 'screenshots/n8n-workflow.png'

os.makedirs(OUT_DIR, exist_ok=True)
prs = Presentation()
# blank layout
blank = prs.slide_layouts[6]

# helper functions

def set_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(blank)
    set_background(slide, (10, 12, 20))
    # Title
    left = Inches(0.6)
    top = Inches(1.0)
    width = Inches(9)
    height = Inches(1.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 255, 255)

    # Subtitle
    sub = slide.shapes.add_textbox(left, Inches(2.3), width, Inches(0.8))
    st = sub.text_frame
    sp = st.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(120, 230, 255)


def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(blank)
    set_background(slide, (8, 10, 16))
    # title
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(9), Inches(0.8)).text_frame
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 240, 255)
    # bullets
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.5)).text_frame
    tx.margin_left = Inches(0)
    for b in bullets:
        p = tx.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(180, 220, 255)


def add_architecture_slide():
    slide = prs.slides.add_slide(blank)
    set_background(slide, (6, 8, 14))
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6)).text_frame
    title.paragraphs[0].text = 'Arquitectura'
    title.paragraphs[0].font.size = Pt(28)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = RGBColor(200,240,255)

    # Draw simple workflow boxes
    lefts = [Inches(0.8), Inches(3.8), Inches(6.8)]
    labels = ['Schedule\nTrigger', 'Code Node\n(JavaScript)', 'Provider\n(Simulated)']
    for i, L in enumerate(lefts):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, Inches(2.0), Inches(2.0), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(10, 200, 230)
        shape.line.color.rgb = RGBColor(0,160,200)
        tf = shape.text_frame
        tf.text = labels[i]
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(6,6,6)

    # Downstream: Google Sheets and Telegram
    g = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(4.0), Inches(3.4), Inches(1.0))
    g.fill.solid(); g.fill.fore_color.rgb = RGBColor(40,180,255)
    g.text_frame.text = 'Google Sheets\n(append rows)'
    g.text_frame.paragraphs[0].font.size = Pt(12)

    t = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(4.0), Inches(3.0), Inches(1.0))
    t.fill.solid(); t.fill.fore_color.rgb = RGBColor(40,180,255)
    t.text_frame.text = 'Telegram\n(alerts)'
    t.text_frame.paragraphs[0].font.size = Pt(12)


def add_screenshot_slide():
    slide = prs.slides.add_slide(blank)
    set_background(slide, (4, 6, 12))
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.5)).text_frame
    title.paragraphs[0].text = 'Workflow Preview'
    title.paragraphs[0].font.size = Pt(26)
    title.paragraphs[0].font.color.rgb = RGBColor(200,240,255)

    if os.path.exists(SCREENSHOT):
        # fit image into slide area
        img = Image.open(SCREENSHOT)
        max_w = Inches(9)
        max_h = Inches(5)
        # add picture
        slide.shapes.add_picture(SCREENSHOT, Inches(0.5), Inches(1.2), width=max_w)
    else:
        # placeholder
        ph = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        p = ph.text_frame.paragraphs[0]
        p.text = 'No screenshot found. Place `screenshots/n8n-workflow.png` to include it.'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(180,200,230)


# Build presentation
add_title_slide('SmartFare Tracker', 'n8n automation — Flight price monitoring')
add_bullet_slide('Resumen', [
    'Monitorea precios de vuelos con n8n',
    'Registra resultados en Google Sheets',
    'Envía alertas por Telegram para precios baratos',
    'Tecnologías: JavaScript, OAuth2, Docker, GitHub'
])
add_architecture_slide()
add_screenshot_slide()
add_bullet_slide('Detalles técnicos', [
    'Ejecución programada (Schedule Trigger)',
    'Nodo de código en JavaScript genera precios simulados',
    'Google Sheets: append rows via OAuth2',
    'Telegram: envíos mediante Bot API'
])
add_bullet_slide('Siguientes pasos', [
    'Integrar proveedor real de vuelos',
    'Mejorar reglas de alerta',
    'Desplegar runners y monitoreo'
])

prs.save(OUT_FILE)
print('Saved', OUT_FILE)
